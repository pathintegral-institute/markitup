import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any
from operator import attrgetter

from ._html_converter import HtmlConverter
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._stream_info import StreamInfo
import pptx


ACCEPTED_MAGIC_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]

ACCEPTED_FILE_CATEGORY = [".pptx"]


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self):
        super().__init__()
        self._html_converter = HtmlConverter()

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    alt_text = ""

                    # Also grab any description embedded in the deck
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        # Unable to get alt text
                        pass

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # If keep_data_uris is True, use base64 encoding for images

                    blob = shape.image.blob
                    content_type = shape.image.content_type or "image/png"
                    b64_string = base64.b64encode(blob).decode("utf-8")
                    md_content += f"\n![{alt_text}](data:{content_type};base64,{b64_string})\n"


                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(shape.shapes, key=attrgetter("top", "left"))
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(slide.shapes, key=attrgetter("top", "left"))
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
